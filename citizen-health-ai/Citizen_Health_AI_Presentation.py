#!/usr/bin/env python3
"""
Generate a professional PowerPoint presentation for Citizen Health AI project.
Run: python Citizen_Health_AI_Presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
DARK_BLUE = RGBColor(31, 78, 121)
LIGHT_BLUE = RGBColor(79, 129, 189)
ACCENT_GREEN = RGBColor(155, 187, 89)
WHITE = RGBColor(255, 255, 255)
DARK_TEXT = RGBColor(51, 51, 51)


def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(60)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.word_wrap = True
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = ACCENT_GREEN
    subtitle_frame.word_wrap = True
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.4))
    footer_frame = footer_box.text_frame
    footer_frame.text = "AI-Powered Health Assistance for Citizens & Field Workers"
    footer_frame.paragraphs[0].font.size = Pt(14)
    footer_frame.paragraphs[0].font.color.rgb = LIGHT_BLUE


def add_content_slide(prs, title, content_points):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.color.rgb = DARK_BLUE
    
    # Title text
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.margin_left = Inches(0.5)
    title_frame.margin_top = Inches(0.15)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(8.5), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = point
        p.font.size = Pt(24)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(12)
        p.space_after = Pt(12)
        p.level = 0


def add_two_column_slide(prs, title, left_title, left_points, right_title, right_points):
    """Add slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.color.rgb = DARK_BLUE
    
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.margin_left = Inches(0.5)
    
    # Left column header
    left_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.2), Inches(0.4))
    left_h_frame = left_header.text_frame
    left_h_frame.text = left_title
    left_h_frame.paragraphs[0].font.size = Pt(20)
    left_h_frame.paragraphs[0].font.bold = True
    left_h_frame.paragraphs[0].font.color.rgb = LIGHT_BLUE
    
    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(4.2), Inches(5.2))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, point in enumerate(left_points):
        if i > 0:
            left_frame.add_paragraph()
        p = left_frame.paragraphs[i]
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(8)
        p.space_after = Pt(8)
    
    # Right column header
    right_header = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.2), Inches(0.4))
    right_h_frame = right_header.text_frame
    right_h_frame.text = right_title
    right_h_frame.paragraphs[0].font.size = Pt(20)
    right_h_frame.paragraphs[0].font.bold = True
    right_h_frame.paragraphs[0].font.color.rgb = LIGHT_BLUE
    
    # Right column content
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.7), Inches(4.2), Inches(5.2))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, point in enumerate(right_points):
        if i > 0:
            right_frame.add_paragraph()
        p = right_frame.paragraphs[i]
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(8)
        p.space_after = Pt(8)


def add_closing_slide(prs, title, text):
    """Add closing slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.word_wrap = True
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    content_frame = content_box.text_frame
    content_frame.text = text
    content_frame.paragraphs[0].font.size = Pt(24)
    content_frame.paragraphs[0].font.color.rgb = ACCENT_GREEN
    content_frame.word_wrap = True


# Slide 1: Title Slide
add_title_slide(prs, "Citizen Health AI", "Digital Health Solutions for India")

# Slide 2: Overview
add_content_slide(prs, "Project Overview", [
    "🏥 AI-powered multilingual health assistant platform",
    "📱 Serves citizens, ASHA workers, and health surveillance officers",
    "🗣️ Voice & text interfaces in Tamil, Kannada, and English",
    "🚀 Deployed as FastAPI microservice with IVR capabilities",
    "💡 Enables accessible healthcare information and support"
])

# Slide 3: Core Features
add_two_column_slide(prs, "Core Features",
    "Citizen Module", [
        "• Portal widget for health queries",
        "• IVR voice interface (WAV)",
        "• Hospital & insurance info",
        "• Vaccination guidance",
        "• Maternal health support"
    ],
    "Field Worker Module", [
        "• ASHA & PHC nurse support",
        "• Protocol & referral lookup",
        "• Patient record updates",
        "• Voice-driven documentation",
        "• Real-time guidance"
    ]
)

# Slide 4: Disease Surveillance
add_content_slide(prs, "Disease Surveillance", [
    "🔍 Outbreak detection & analysis for districts",
    "📊 Statistical anomaly detection (z-score analysis)",
    "🚨 Priority disease scanning across regions",
    "📈 Weekly trend analysis & alerts",
    "⚠️ Pre-alert & critical alert classification"
])

# Slide 5: Technical Architecture
add_two_column_slide(prs, "Technical Stack",
    "Backend & API", [
        "• FastAPI framework",
        "• Python 3.x",
        "• Uvicorn ASGI server",
        "• CORS enabled",
        "• RESTful endpoints"
    ],
    "AI & Audio", [
        "• Sarvam AI integration",
        "• Speech-to-Text (Saarika)",
        "• Text-to-Speech (Bulbul)",
        "• Language detection",
        "• OpenAI LLM"
    ]
)

# Slide 6: API Endpoints
add_content_slide(prs, "REST API Endpoints", [
    "POST /citizen/chat - Text-based health queries",
    "POST /citizen/voice - Voice IVR interaction",
    "POST /worker/chat - Worker protocol queries",
    "POST /worker/voice - Worker voice input",
    "POST /worker/voice-update - Patient record updates",
    "POST /surveillance/analyse - Single district analysis",
    "POST /surveillance/scan - Multi-disease scanning"
])

# Slide 7: Multilingual Support
add_content_slide(prs, "Multilingual Capabilities", [
    "🇮🇳 Tamil (Tamil Nadu & Puducherry)",
    "🇮🇳 Kannada (Karnataka)",
    "🇬🇧 English (National access)",
    "🎯 Automatic language detection",
    "📝 Context-aware translations",
    "🔊 Native speaker voice synthesis"
])

# Slide 8: IVR System
add_two_column_slide(prs, "Interactive Voice Response (IVR)",
    "Session Management", [
        "• Session creation & tracking",
        "• Persistent conversation history",
        "• Multi-turn dialogue support",
        "• Language switching",
        "• Session timeout handling"
    ],
    "Voice Pipeline", [
        "• WAV/WebM audio input",
        "• Saarika speech-to-text",
        "• LLM inference",
        "• Bulbul text-to-speech",
        "• Base64 audio response"
    ]
)

# Slide 9: Data & Models
add_content_slide(prs, "Data & AI Models", [
    "📊 NHIM knowledge base integration (citizen module)",
    "🏥 Patient record structuring for HMIS",
    "📈 Statistical outbreak detection algorithms",
    "🧠 Large Language Model (LLM) for Q&A",
    "🔐 Secure patient data handling"
])

# Slide 10: Use Cases
add_two_column_slide(prs, "Use Cases",
    "Citizens", [
        "• Find nearest hospitals",
        "• Learn about insurance",
        "• Check vaccination schedule",
        "• Access maternal health info",
        "• Voice-enabled accessibility"
    ],
    "Healthcare Workers", [
        "• Quick protocol reference",
        "• Patient data entry",
        "• Referral coordination",
        "• Field-to-system sync",
        "• Worker support 24/7"
    ]
)

# Slide 11: Impact
add_content_slide(prs, "Expected Impact", [
    "✓ Improved health literacy among citizens",
    "✓ Faster decision-making for field workers",
    "✓ Early disease outbreak detection",
    "✓ Reduced healthcare access barriers",
    "✓ Scalable to 1M+ users",
    "✓ Cost-effective health intervention"
])

# Slide 12: Deployment & Scalability
add_two_column_slide(prs, "Deployment & Scalability",
    "Infrastructure", [
        "• Cloud-agnostic FastAPI",
        "• Docker containerization",
        "• Load balancing ready",
        "• Horizontal scaling",
        "• CDN for static assets"
    ],
    "Performance", [
        "• ~2-5s voice round-trip",
        "• Sub-second text response",
        "• Graceful TTS degradation",
        "• Error handling & logging",
        "• Health check endpoints"
    ]
)

# Slide 13: Future Roadmap
add_content_slide(prs, "Future Enhancements", [
    "📲 Mobile app (iOS/Android)",
    "🔗 Integration with HMIS systems",
    "📍 Geolocation-based services",
    "🤖 Advanced analytics & dashboards",
    "🌐 Additional language support",
    "💬 Chatbot for web & social media"
])

# Slide 14: Team & Collaboration
add_content_slide(prs, "Project Team", [
    "👨‍💼 Project Lead: Atul",
    "👥 Multi-disciplinary team:",
    "  • AI/ML Engineers",
    "  • Backend Developers",
    "  • Health Domain Experts",
    "🤝 Collaboration with healthcare stakeholders",
    "📚 Open to contributions & feedback"
])

# Slide 15: Conclusion
add_closing_slide(prs, "Building Inclusive Digital Health",
    "Empowering every citizen with AI-driven health guidance, one conversation at a time.")

# Save presentation
output_file = "Citizen_Health_AI_Presentation.pptx"
prs.save(output_file)
print(f"✅ Presentation created successfully: {output_file}")
